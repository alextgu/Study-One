import json
import logging
import re
import os
import asyncio
from pathlib import Path

from fastapi import Depends, FastAPI, HTTPException, UploadFile, File
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel, field_validator
from services import GeminiService, get_supabase
from services.gamification import (
    XP_FLASHCARD_SESSION,
    XP_QUIZ_PERFECT_BONUS as PERFECT_SCORE_BONUS,
    XP_QUIZ_PER_CORRECT as XP_CORRECT,
    award_flashcard_session_xp,
    quiz_attempt_xp_breakdown,
    record_quiz_attempt_xp,
)
from services.badge_trigger import evaluate_and_award as evaluate_badge_triggers
from middleware.auth import require_user, UserPayload, user_for_generate
from typing import List, Optional
from enum import Enum
from datetime import datetime, timezone
from pypdf import PdfReader
from pptx import Presentation

logger = logging.getLogger(__name__)

from prompts.study_gen_v1 import build_study_generation_prompt, validate_quiz_quality, build_summary_only_prompt
from prompts.flashcard_gen_v1 import build_flashcard_generation_prompt
from prompts.quiz_gen_v1 import build_quiz_generation_prompt
from prompts.slide_text_gen_v1 import build_chunk_summary_prompt, build_merge_outline_prompt

app = FastAPI(title="Socrato")
UPLOAD_DIR = Path(__file__).resolve().parent / "uploads" / "slides"
MAX_UPLOAD_SIZE_BYTES = 25 * 1024 * 1024
CHUNKING_THRESHOLD_CHARS = 18000

# Initialize Gemini service
gemini_service = GeminiService()

# Enable CORS for frontend communication
app.add_middleware(
    CORSMiddleware,
    allow_origins=["http://localhost:3000"],  # Next.js dev server
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

def check_empty_text(v: str) -> str:
    if not v or not v.strip():
        raise ValueError("text must not be empty")
    return v


# ============================================
# REQUEST/RESPONSE SCHEMAS
# Mirrors shared/types.ts contract
# ============================================

class GenerateRequest(BaseModel):
    """
    Request body for POST /api/v1/generate
    - text: The user's study notes to process
    """
    text: str

    @field_validator("text")
    @classmethod
    def text_must_not_be_empty(cls, v: str) -> str:
        return check_empty_text(v)


class StudyPackRequest(GenerateRequest):
    """
    Request body for POST /generate-study-pack
    - text: The user's study notes to process
    """
    @field_validator("text")
    @classmethod
    def text_length_constraint(cls, v: str) -> str:
        v = check_empty_text(v)
        stripped = v.strip()
        # validate length
        if len(stripped) < 10:
            raise ValueError(f"text must not be less than 10 characters")
        if len(stripped) > 10000:
            raise ValueError("text must not be more than 10000 characters")
        return v



class QuizQuestion(BaseModel):
    """A single quiz question with multiple choice options"""
    question: str
    options: List[str]
    answer: str
    topic: str
    # Simple one-paragraph explanation of why the correct answer is right
    # and why the other options are wrong.
    correctionExplanation: Optional[str] = None


class GenerateResponse(BaseModel):
    """
    Response from POST /api/v1/generate
    - summary: Array of bullet point summaries
    - quiz: Array of quiz questions with options and answers
    """
    summary: List[str]
    quiz: List[QuizQuestion]


class AnkiRating(str, Enum):
    AGAIN = "again"
    HARD = "hard"
    GOOD = "good"
    EASY = "easy"

class FlashcardReviewRequest(BaseModel):
    flashcard_set_id: str
    card_index: int
    rating: AnkiRating

class FlashcardReviewResponse(BaseModel):
    xp_awarded: int
    total_xp: int
    already_reviewed: bool


class QuizResultRequest(BaseModel):
    """Request body for POST /api/v1/quiz/result - submit quiz completion for XP."""
    correct: int
    total: int
    quiz_id: Optional[str] = None


class XpAwardLine(BaseModel):
    """One line item explaining part of an XP award (UI + analytics)."""
    reason: str
    xp: int


class QuizResultResponse(BaseModel):
    """Deprecated: XP is recorded on POST /api/v1/quiz/attempt. Kept for API compatibility."""
    applied: bool
    xp_awarded: int
    user_stats: dict
    xp_breakdown: List[XpAwardLine] = []
    notice: Optional[str] = None


class FlashcardSessionCompleteRequest(BaseModel):
    """Request body for POST /api/v1/flashcards/session-complete."""
    flashcard_set_id: str


class FlashcardSessionCompleteResponse(BaseModel):
    """Response from flashcard session completion."""
    applied: bool
    xp_awarded: int
    user_stats: dict
    xp_breakdown: List[XpAwardLine] = []

class QuizExplanationRequest(BaseModel):
    """
    Request body for POST /api/v1/quiz/explain
    - question/options/answer: the MC question context
    - user_answer: the option the student chose (optional)
    - correction_explanation: any existing explanation shown to the user (optional)
    - followup_prompt: optional follow-up question from the student
    """
    question: str
    options: List[str]
    answer: str
    user_answer: Optional[str] = None
    correction_explanation: Optional[str] = None
    followup_prompt: Optional[str] = None


class QuizExplanationResponse(BaseModel):
    """Plain-text explanation for a quiz question."""
    explanation: str

# ============================================
# STUDY PACK HELPER FUNCTIONS
# ============================================
def clean_response(response):
    """
    Clean up Gemini response by removing markdown code blocks
    """
    # Clean up response if it has markdown code blocks
    cleaned = response.strip() 
    # remove opening markdown code fence
    cleaned = re.sub(r'^```[a-z]*\n?', '', cleaned) 
    # remove closing markdown code fence
    cleaned = re.sub(r'```$', '', cleaned)

    return cleaned.strip()


def validate_data(data):
    """
    Validate the study pack has all the required fields and return the list of quiz questions
    """
    # Validate required fields exist
    if not isinstance(data.get("summary"), list):
        raise ValueError("Response missing 'summary' array")
    if not isinstance(data.get("quiz"), list):
        raise ValueError("Response missing 'quiz' array")
    
    # Parse quiz questions with validation
    quiz_questions = []
    for i, q in enumerate(data.get("quiz", [])):
        if not isinstance(q, dict):
            raise ValueError(f"Quiz item {i} is not an object")
        if "question" not in q:
            raise ValueError(f"Quiz item {i} missing 'question' field")
        if "options" not in q or not isinstance(q["options"], list):
            raise ValueError(f"Quiz item {i} missing 'options' array")
        if "answer" not in q:
            raise ValueError(f"Quiz item {i} missing 'answer' field")
        
        quiz_questions.append(QuizQuestion(
            question=q["question"],
            options=q["options"],
            answer=q["answer"],
            topic=(q.get("topic") or "").strip() or "General",
            correctionExplanation=q.get("correctionExplanation")
        ))

    return quiz_questions



# ============================================
# ROUTES
# ============================================

@app.get("/")
def root():
    return {}


@app.get("/health")
def check_health():
    return {"status": "ok"}


@app.get("/api/v1/me")
async def get_current_user(user: UserPayload = Depends(require_user)):
    """Return the authenticated user's identity. 401 if not logged in."""
    return {
        "user_id": user["user_id"],
        "email": user.get("email"),
        "role": user.get("role"),
    }


@app.post("/api/v1/generate", response_model=GenerateResponse)
async def generate_study_materials(
    request: GenerateRequest,
    _user: Optional[UserPayload] = Depends(user_for_generate),
):
    """Generate study materials from user notes. Auth controlled by REQUIRE_AUTH_FOR_GENERATE."""

    # Build prompt using the centralized prompt system
    prompt = build_study_generation_prompt(
        user_notes=request.text,
        include_examples=True  # Include few-shot examples for better quality
    )

    response = await gemini_service.call_gemini(prompt)
    
    if response is None:
        raise HTTPException(
            status_code=500,
            detail="Failed to generate study materials. Please try again."
        )
    
    # Parse the JSON response from Gemini
    try:
        # Clean up response if it has markdown code blocks
        cleaned = clean_response(response)
        
        data = json.loads(cleaned)
        
        quiz_questions = validate_data(data)
        
        # Optional: Run quality checks on the quiz
        quality_warnings = validate_quiz_quality(data.get("quiz", []))
        if quality_warnings:
            # print(f"[generate] Quality warnings: {quality_warnings}")
            # Can log these or return them to the frontend in the future
            logger.info("Quiz quality warnings count: %d", len(quality_warnings))


        return GenerateResponse(
            summary=data.get("summary", []),
            quiz=quiz_questions
        )
    except json.JSONDecodeError as e:
        logger.warning("Failed to parse Gemini JSON: %s", e)
        logger.debug("Raw Gemini response length: %s", len(response) if response else 0)
        raise HTTPException(
            status_code=500,
            detail="Failed to parse AI response as JSON. Please try again."
        )
    except (KeyError, TypeError, ValueError) as e:
        logger.warning("Invalid Gemini response structure: %s", e)
        logger.debug("Raw Gemini response length: %s", len(response) if response else 0)
        raise HTTPException(
            status_code=500,
            detail=f"Invalid AI response format: {str(e)}"
        )


# ============================================
# STUDY PACK ROUTE
# ============================================


@app.post("/generate-study-pack", response_model=GenerateResponse)
async def generate_study_pack(
    request: StudyPackRequest,
    _user: Optional[UserPayload] = Depends(user_for_generate),
):
    """Generate a study pack from user notes. Auth controlled by REQUIRE_AUTH_FOR_GENERATE."""
    prompt = build_study_generation_prompt(
        user_notes=request.text,
        include_examples=True,
    )
           
    # Call Gemini API
    response = await gemini_service.call_gemini(prompt)
    
    if response is None:
        raise HTTPException(
            status_code=500,
            detail="Gemini unavailable. Please try again."
        )
        
    try:
        # Clean up response if it has markdown code blocks
        cleaned = clean_response(response)
        
        data = json.loads(cleaned)
        
        # Validate required fields exist
        quiz_questions = validate_data(data)

        quality_warnings = validate_quiz_quality(data.get("quiz", []))
        if quality_warnings:
            logger.info("Quiz quality warnings count: %d", len(quality_warnings))
        
        return GenerateResponse(
            summary=data['summary'],
            quiz=quiz_questions
        )
    
    except json.JSONDecodeError as e:
        logger.warning("Failed to parse Gemini JSON: %s", e)
        logger.debug("Raw Gemini response: %s", response)
        raise HTTPException(
            status_code=500,
            detail="Failed to parse AI response as JSON. Please try again."
        )
    except (KeyError, TypeError, ValueError) as e:
        logger.warning("Invalid Gemini response structure: %s", e)
        logger.debug("Raw Gemini response: %s", response)
        raise HTTPException(
            status_code=500,
            detail=f"Invalid AI response format: {str(e)}"
        )
    

# ============================================
# QUIZ SCHEMA
# ============================================

class GenerateQuizResponse(BaseModel):
    """
    Response from POST /api/v1/quiz
    - quiz: Array of quiz questions with options, answers, and a linked topic
    """
    quiz_set_id: str
    quiz: list[QuizQuestion]

class SlideQuizRegenerateRequest(BaseModel):
    """Request body for POST /api/v1/slides/quiz/regenerate."""
    text: str

    @field_validator("text")
    @classmethod
    def text_must_not_be_empty(cls, v: str) -> str:
        v = check_empty_text(v)
        if len(v.strip()) < 10:
            raise ValueError("text must not be less than 10 characters")
        return v

class QuestionAnswer(BaseModel):
    question_index: int
    selected_answer: str

    @field_validator("selected_answer")
    @classmethod
    def answer_not_empty(cls, v: str) -> str:
        v = check_empty_text(v)
        return v

class QuizSubmitRequest(BaseModel):
    quiz_id: str
    answers: list[QuestionAnswer]

class QuestionResult(BaseModel):
    question_index: int
    question: str
    selected_answer: str
    correct_answer: str
    is_correct: bool
    topic: str
    correction_explanation: Optional[str] = None

class QuizSubmitResponse(BaseModel):
    attempt_id: str
    quiz_set_id: str
    score: float
    total_correct: int
    total_questions: int
    xp_awarded: int
    xp_breakdown: List[XpAwardLine] = []
    results: list[QuestionResult]


def calc_xp(correct: int, total: int) -> int:
    xp, _lines = quiz_attempt_xp_breakdown(correct=correct, total=total)
    return xp


# ============================================
# QUIZ HELPER
# ============================================

def parse_and_validate_quiz(raw_response: str) -> list[QuizQuestion]:
    """Parse Gemini JSON and validate structure for quiz.""" 
    cleaned = clean_response(raw_response)
    try:
        data = json.loads(cleaned)
    
        raw_quiz = data.get("quiz")
        if not isinstance(raw_quiz, list):
            raise ValueError("Response missing 'quiz' array")
        # check the number of quiz generated
        if len(raw_quiz) < 5 or len(raw_quiz) > 10:
            raise ValueError(f"Expected 5-10 questions, got {len(raw_quiz)}.")

        # Parse quiz questions with validation
        quiz_questions = []
        for i, q in enumerate(data.get("quiz", [])):
            if not isinstance(q, dict):
                raise ValueError(f"Quiz item {i} is not an object")
            if "question" not in q:
                raise ValueError(f"Quiz item {i} missing 'question' field")
            if "options" not in q or not isinstance(q["options"], list):
                raise ValueError(f"Quiz item {i} missing 'options' array")
            if "answer" not in q:
                raise ValueError(f"Quiz item {i} missing 'answer' field")
            if "topic" not in q or not q["topic"].strip():
                raise ValueError(f"Quiz item {i} missing 'topic' field")
            if q["answer"] not in q["options"]:
                raise ValueError(f"Quiz item {i} 'answer' not in 'options'")
            

            quiz_questions.append(QuizQuestion(
                question=q["question"],
                options=q["options"],
                answer=q["answer"],
                topic=q["topic"],
                correctionExplanation=q.get("correctionExplanation"),
            ))
        
        return quiz_questions
    except json.JSONDecodeError as e:
        logger.warning("[quiz] Failed to parse Gemini JSON: %s", e)
        logger.debug("[quiz] Raw Gemini response: %s", raw_response)
        raise HTTPException(
            status_code=500,
            detail="Failed to parse AI response as JSON. Please try again."
        )
    except (KeyError, TypeError, ValueError) as e:
        logger.warning("[quiz] Invalid Gemini response structure: %s", e)
        logger.debug("[quiz] Raw Gemini response: %s", raw_response)
        raise HTTPException(
            status_code=500,
            detail=f"Invalid AI response format: {str(e)}"
        )

def grade_quiz(answers: list[QuestionAnswer], questions: list[QuizQuestion]) -> list[QuestionResult]:
    # get all the answers from the user
    user_answers = {a.question_index: a.selected_answer for a in answers}

    quiz_results = []

    for i, q in enumerate(questions):
        ans = user_answers.get(i, "")
        is_correct = ans == q.answer

        result = QuestionResult(
            question_index=i,
            question=q.question,
            selected_answer=ans,
            correct_answer=q.answer,
            is_correct=is_correct,
            topic=q.topic,
            correction_explanation=q.correctionExplanation,
        ) 
        quiz_results.append(result)
    
    return quiz_results

def validate_submit_quiz_request(request: QuizSubmitRequest) -> None:
    try: 
        if not request.quiz_id:
            raise ValueError("Missing 'quiz_id' field.")
        # Allow an empty list to flow through to the length check later;
        # we only treat a completely absent answers field as invalid here.
        if request.answers is None:  # type: ignore[comparison-overlap]
            raise ValueError("Missing 'answers' field.")
        if not isinstance(request.quiz_id, str):
            raise ValueError(f"Invalid 'quiz_id' field.")  
        if not isinstance(request.answers, list):
            raise ValueError(f"Invalid 'answers' field.")

        for a in request.answers:
            if not isinstance(a, QuestionAnswer): 
                raise ValueError(f"Invalid 'answers' field.")
            if not isinstance(a.question_index, int):
                raise ValueError(f"Invalid 'question_index' field.")
            if not isinstance(a.selected_answer, str):
                raise ValueError(f"Invalid 'selected_answer' field.")
            if a.selected_answer.strip() == "":
                raise ValueError(f"Invalid 'selected_answer' field.")
    except (KeyError, TypeError, ValueError) as e:
        logger.warning("[quiz submit] Invalid Request: %s", e)
        raise HTTPException(
            status_code=422,
            detail=f"Invalid request format: {str(e)}"
        )

@app.post("/api/v1/quiz", response_model=GenerateQuizResponse)
async def generate_quiz_questions(
    request: StudyPackRequest,
    user: UserPayload | None = Depends(user_for_generate),
):
    """Generate MC Quiz from user notes. Store quiz in supabase."""

    prompt = build_quiz_generation_prompt(content=request.text)
    response = await gemini_service.call_gemini(prompt)
    
    if response is None:
        raise HTTPException(
            status_code=500,
            detail="Failed to generate quiz. Please try again."
        )
    
    # clean up and parse the raw response from Gemini
    quiz_questions = parse_and_validate_quiz(response)

    # store quiz into supabase
    sb = get_supabase()
    try:
        result = sb.table("quiz").insert({
            "user_id": user["user_id"] if user else "00000000-0000-0000-0000-000000000001",
            "source_text": request.text,
            "questions": [q.model_dump() for q in quiz_questions],
        }).execute()
        quiz_set_id = result.data[0]["id"]
    except Exception as e:
        print(f"[quiz] DB insert failed: {e}")
        raise HTTPException(
            status_code=500,
            detail="Failed to store quiz. Please try again."
        )


    return GenerateQuizResponse(quiz_set_id=quiz_set_id, quiz=quiz_questions)

@app.post("/api/v1/quiz/attempt", response_model=QuizSubmitResponse)
async def submit_quiz_attempt(
    request: QuizSubmitRequest,
    user: UserPayload = Depends(require_user),
):
    """Submit a quiz attempt: grade answers, store the attempt and results, return score and XP."""
    
    # check if the request is valid
    validate_submit_quiz_request(request)
    
    sb = get_supabase()
    user_id = user["user_id"]

    # get the corresponding quiz from the database to get the correct answer
    try: 
        quiz_data = sb.table("quiz") \
        .select("*") \
        .eq("id", request.quiz_id) \
        .single() \
        .execute()

        if not quiz_data:
            raise HTTPException(status_code=404, detail=f"Quiz {request.quiz_id} not found.")

        questions = [QuizQuestion(**q) for q in quiz_data.data["questions"]]

    except HTTPException:
        raise
    except Exception as e:
        logger.warning(f"[quiz submit] DB query failed: {e}")
        raise HTTPException(
            status_code=500,
            detail="Failed to retrieve quiz. Please try again."
        )

     # user not answer all questions
    if len(request.answers) != len(questions):
        raise HTTPException(
            status_code=422,
            detail=f"Expected {len(questions)} answers but received {len(request.answers)}. Please answer all questions before submitting."
        )

    # if user submit duplicate answers for the same question
    submitted_indices = [a.question_index for a in request.answers]
    if sorted(submitted_indices) != list(range(len(questions))):
        raise HTTPException(
            status_code=422,
            detail="Answers must contain exactly one response per question with no duplicates."
        )

   # check user's answers are in question options
    for ans in request.answers:
        if ans.question_index < 0 or ans.question_index >= len(questions):
            raise HTTPException(
                status_code=422,
                detail=f"Invalid 'question_index' {ans.question_index}. Must be between 0 and {len(questions)-1}."
            )
        if ans.selected_answer not in questions[ans.question_index].options:
            raise HTTPException(
                status_code=422,
                detail=f"Question {ans.question_index}: '{ans.selected_answer}' is not a valid option."
            )

    # grade the user's response
    question_result = grade_quiz(request.answers, questions)
    correct = sum(1 for qr in question_result if qr.is_correct)
    total = len(questions)
    score = round((correct / total) * 100, 2)
    xp, breakdown_lines = quiz_attempt_xp_breakdown(correct=correct, total=total)

    # store the user's attempt of the quiz into supabase
    try:
        result = sb.table("quiz_attempt").insert({
            "user_id": user["user_id"] if user else "00000000-0000-0000-0000-000000000001",
            "quiz_set_id": request.quiz_id,
            "score": score, 
            "total_correct": correct, 
            "total_questions": total, 
            "xp_awarded": xp, 
            "results": [qr.model_dump() for qr in question_result],
        }).execute()
        attempt_id = result.data[0]["id"]
    except Exception as e:
        logger.warning(f"[quiz submit] DB insert failed: {e}")
        raise HTTPException(
            status_code=500,
            detail="Failed to store quiz attempt. Please try again."
        )

    activity_md = {
        "quiz_set_id": request.quiz_id,
        "attempt_id": attempt_id,
        "total_correct": correct,
        "total_questions": total,
        "score": score,
        "xp_breakdown": breakdown_lines,
    }
    try:
        rec = record_quiz_attempt_xp(
            user_id=user_id,
            xp_awarded=xp,
            metadata=activity_md,
        )
        us = rec.get("user_stats")
        if xp > 0 and isinstance(us, dict):
            try:
                evaluate_badge_triggers(user_id, user_stats=us)
            except Exception as e:
                logger.warning("Badge trigger evaluation failed after quiz attempt: %s", e)
    except Exception as e:
        logger.warning("apply_quiz_attempt_record failed: %s", e)

    submit_response = QuizSubmitResponse(
        attempt_id=attempt_id,
        quiz_set_id=request.quiz_id,
        score=score,
        total_correct=correct, 
        total_questions=total,
        xp_awarded=xp,
        xp_breakdown=[XpAwardLine(**line) for line in breakdown_lines],
        results=question_result
    )

    return submit_response

@app.post("/api/v1/quiz/explain", response_model=QuizExplanationResponse)
async def explain_quiz_answer(
    request: QuizExplanationRequest,
    _user: Optional[UserPayload] = Depends(user_for_generate),
):
    """
    Return a short, focused explanation for why the correct answer is right
    and/or answer a follow-up question about this specific quiz item.
    """
    base_context = {
        "question": request.question,
        "options": request.options,
        "correct_answer": request.answer,
        "user_answer": request.user_answer,
        "existing_explanation": request.correction_explanation,
    }

    if request.followup_prompt:
        prompt = f"""You are helping a student understand a multiple-choice question.

Question and options (JSON):
{json.dumps(base_context, ensure_ascii=False, indent=2)}

The student has a follow-up question:
\"\"\"{request.followup_prompt.strip()}\"\"\"

Rules:
- Use ONLY the information in the JSON above.
- If the follow-up asks about something beyond that context, say that you don't know.
- Focus on explaining the reasoning clearly and simply.
- Respond in 2–4 short sentences.
- Respond as plain text only (no lists, bullet points, or markdown)."""
    else:
        prompt = f"""You are helping a student understand why their multiple-choice answer was incorrect.

Question and options (JSON):
{json.dumps(base_context, ensure_ascii=False, indent=2)}

Task:
- Briefly explain why the correct answer is right and why the student's answer, if given, is wrong.

Rules:
- Use ONLY the information in the JSON above.
- Do NOT introduce outside facts or extra background knowledge.
- Respond in 2–4 short sentences.
- Respond as plain text only (no lists, bullet points, or markdown)."""

    raw_response = await gemini_service.call_gemini(prompt)

    if raw_response is None:
        raise HTTPException(
            status_code=500,
            detail="Gemini unavailable. Please try again.",
        )

    cleaned = clean_response(raw_response)
    explanation = cleaned.strip()

    if not explanation:
        raise HTTPException(
            status_code=500,
            detail="AI returned an empty explanation. Please try again.",
        )

    return QuizExplanationResponse(explanation=explanation)


@app.post("/api/v1/quiz/result", response_model=QuizResultResponse)
async def submit_quiz_result(
    request: QuizResultRequest,
    user: UserPayload = Depends(require_user),
):
    """
    Deprecated: XP and streaks are updated when you POST /api/v1/quiz/attempt.
    This endpoint returns current stats without awarding duplicate XP.
    """
    if request.total <= 0 or request.correct < 0 or request.correct > request.total:
        raise HTTPException(
            status_code=400,
            detail="Invalid correct/total: total must be > 0 and correct must be in [0, total].",
        )
    sb = get_supabase()
    try:
        row = (
            sb.table("user_stats")
            .select("user_id, xp_total, level, current_streak_days, longest_streak_days, last_active_at")
            .eq("user_id", user["user_id"])
            .maybe_single()
            .execute()
        )
        stats = row.data if row.data else {}
    except Exception as e:
        logger.warning("Quiz result stats fetch failed: %s", e)
        stats = {}
    return QuizResultResponse(
        applied=False,
        xp_awarded=0,
        user_stats=stats,
        xp_breakdown=[],
        notice="Quiz XP is recorded when you submit answers to POST /api/v1/quiz/attempt.",
    )


# ============================================
# FLASHCARD SCHEMAS
# ============================================

class FlashcardRequest(BaseModel):
    text: Optional[str] = None
    topic: Optional[str] = None

    @field_validator("text", "topic")
    @classmethod
    def strip_strings(cls, v: Optional[str]) -> Optional[str]:
        if v is None:
            return v
        v = v.strip()
        return v or None

    def model_post_init(self, __context) -> None:
        if not self.text and not self.topic:
            raise ValueError("Either 'text' or 'topic' must be provided.")


class Flashcard(BaseModel):
    question: str
    answer: str


class FlashcardResponse(BaseModel):
    flashcard_set_id: str
    flashcards: List[Flashcard]


# ============================================
# FLASHCARD HELPERS
# ============================================

def parse_and_validate_flashcards(raw_response: str) -> List[Flashcard]:
    """Parse Gemini JSON and validate structure for flashcards."""
    # Use the existing cleaning logic from clean_response
    cleaned = clean_response(raw_response)

    try:
        data = json.loads(cleaned)
    except json.JSONDecodeError as e:
        print(f"[flashcards] Failed to parse JSON: {e}")
        print(f"[flashcards] Raw response: {raw_response}")
        raise HTTPException(
            status_code=500,
            detail="Failed to parse AI response as JSON. Please try again."
        )

    flashcards_raw = data.get("flashcards")
    if not isinstance(flashcards_raw, list):
        raise HTTPException(
            status_code=500,
            detail="Invalid AI response: 'flashcards' must be an array."
        )
    if len(flashcards_raw) != 10:
        raise HTTPException(
            status_code=500,
            detail=f"Invalid AI response: expected 10 flashcards, got {len(flashcards_raw)}."
        )

    flashcards = []
    for i, fc in enumerate(flashcards_raw):
        if not isinstance(fc, dict):
            raise HTTPException(status_code=500, detail=f"Flashcard {i} is not an object.")
        q = fc.get("question")
        a = fc.get("answer")
        if not isinstance(q, str) or not q.strip():
            raise HTTPException(status_code=500, detail=f"Flashcard {i} missing valid 'question'.")
        if not isinstance(a, str) or not a.strip():
            raise HTTPException(status_code=500, detail=f"Flashcard {i} missing valid 'answer'.")
        flashcards.append(Flashcard(question=q.strip(), answer=a.strip()))

    return flashcards

@app.post("/api/v1/flashcards", response_model=FlashcardResponse)
async def generate_flashcards(
    request: FlashcardRequest, 
    user: Optional[UserPayload] = Depends(user_for_generate)):
    """
    Generate 10 Q/A flashcards from notes or a topic, store in Supabase.
    """
    content = request.text if request.text else request.topic
    mode = "notes" if request.text else "topic"

    prompt = build_flashcard_generation_prompt(
        content=content,
        mode=mode,
    )

    response = await gemini_service.call_gemini(prompt)
    if response is None:
        raise HTTPException(
            status_code=500,
            detail="Failed to generate flashcards. Please try again."
        )

    flashcards = parse_and_validate_flashcards(response)

    # Store in Supabase
    sb = get_supabase()
    try:
        result = sb.table("flashcards").insert({
            "user_id": user["user_id"] if user else "00000000-0000-0000-0000-000000000001",
            "source_text": request.text,
            "topic": request.topic,
            "cards": [fc.model_dump() for fc in flashcards],
        }).execute()
        flashcard_set_id = result.data[0]["id"]
    except Exception as e:
        print(f"[flashcards] DB insert failed: {e}")
        raise HTTPException(
            status_code=500,
            detail="Failed to store flashcards. Please try again."
        )

    return FlashcardResponse(flashcard_set_id=flashcard_set_id, flashcards=flashcards)


@app.post("/api/v1/flashcards/session-complete", response_model=FlashcardSessionCompleteResponse)
async def complete_flashcard_session(
    request: FlashcardSessionCompleteRequest,
    user: UserPayload = Depends(require_user),
):
    """Record flashcard session completion for XP. Awards 10 XP. Idempotent per day."""
    try:
        result = award_flashcard_session_xp(
            user_id=user["user_id"],
            session_id=request.flashcard_set_id,
        )
    except RuntimeError as e:
        logger.warning("Flashcard session apply_activity failed: %s", e)
        raise HTTPException(status_code=500, detail="Failed to record session.")
    try:
        evaluate_badge_triggers(user["user_id"], user_stats=result["user_stats"])
    except Exception as e:
        logger.warning("Badge trigger evaluation failed after flashcard session: %s", e)
    fc_breakdown: List[XpAwardLine] = []
    if result["applied"] and result["xp_awarded"] > 0:
        fc_breakdown = [
            XpAwardLine(reason="Flashcard session completed", xp=XP_FLASHCARD_SESSION),
        ]
    return FlashcardSessionCompleteResponse(
        applied=result["applied"],
        xp_awarded=result["xp_awarded"],
        user_stats=result["user_stats"],
        xp_breakdown=fc_breakdown,
    )


@app.post("/api/v1/flashcards/review", response_model=FlashcardReviewResponse)
async def submit_flashcard_review(
    request: FlashcardReviewRequest,
    user: UserPayload = Depends(require_user),
):
    """Record an Anki-style flashcard review. XP is awarded via POST /api/v1/flashcards/session-complete."""
    sb = get_supabase()
    user_id = user["user_id"]
    today = datetime.now(timezone.utc).date().isoformat()

    # Check for duplicate review today
    existing = sb.table("flashcard_reviews") \
        .select("id") \
        .eq("user_id", user_id) \
        .eq("flashcard_set_id", request.flashcard_set_id) \
        .eq("card_index", request.card_index) \
        .gte("reviewed_at", today) \
        .execute()

    if existing.data:
        stats_result = sb.table("user_stats").select("xp_total").eq("user_id", user_id).maybe_single().execute()
        total_xp = stats_result.data["xp_total"] if stats_result.data else 0
        return FlashcardReviewResponse(xp_awarded=0, total_xp=total_xp, already_reviewed=True)

    try:
        sb.table("flashcard_reviews").insert({
            "user_id": user_id,
            "flashcard_set_id": request.flashcard_set_id,
            "card_index": request.card_index,
            "rating": request.rating.value,
            "xp_awarded": 0,
        }).execute()
    except Exception as e:
        logger.warning("Failed to insert flashcard_reviews: %s", e)
        raise HTTPException(status_code=500, detail="Failed to record review. Please try again.")

    stats_result = sb.table("user_stats").select("xp_total").eq("user_id", user_id).maybe_single().execute()
    total_xp = stats_result.data["xp_total"] if stats_result.data else 0

    return FlashcardReviewResponse(xp_awarded=0, total_xp=total_xp, already_reviewed=False)


@app.get("/api/v1/flashcards/{flashcard_set_id}/session-summary")
async def get_session_summary(
    flashcard_set_id: str,
    user: Optional[UserPayload] = Depends(user_for_generate),
):
    """Return today's ratings for this flashcard set."""
    sb = get_supabase()
    user_id = user["user_id"] if user else "00000000-0000-0000-0000-000000000001"
    today = datetime.now(timezone.utc).date().isoformat()

    result = sb.table("flashcard_reviews") \
        .select("card_index, rating, reviewed_at") \
        .eq("user_id", user_id) \
        .eq("flashcard_set_id", flashcard_set_id) \
        .gte("reviewed_at", today) \
        .order("reviewed_at", desc=True) \
        .execute()

    return {"reviews": result.data}


RATING_PRIORITY = {"again": 0, "hard": 1, "good": 2, "easy": 3}

@app.get("/api/v1/flashcards/{flashcard_set_id}/history")
async def get_card_history(
    flashcard_set_id: str,
    user: Optional[UserPayload] = Depends(user_for_generate),
):
    """Return most recent rating per card, sorted by again -> hard -> good -> easy."""
    sb = get_supabase()
    user_id = user["user_id"] if user else "00000000-0000-0000-0000-000000000001"

    result = sb.table("flashcard_reviews") \
        .select("card_index, rating, reviewed_at") \
        .eq("user_id", user_id) \
        .eq("flashcard_set_id", flashcard_set_id) \
        .order("reviewed_at", desc=True) \
        .execute()

    # Keep only most recent rating per card
    seen = set()
    latest_per_card = []
    for row in result.data:
        if row["card_index"] not in seen:
            seen.add(row["card_index"])
            latest_per_card.append(row)

    # Sort so again/hard come first
    latest_per_card.sort(key=lambda x: RATING_PRIORITY.get(x["rating"], 99))

    return {"history": latest_per_card}


# SLIDE UPLOAD SYSTEM
class SlideStudyPackResponse(BaseModel):
    """Response from POST /api/v1/slides/study-pack."""
    file_name: str
    stored_path: str
    quiz_set_id: str
    flashcard_set_id: str
    extracted_text: str
    summary: List[str]
    quiz: List[QuizQuestion]
    flashcards: List[Flashcard]

def extract_pdf(file: Path) -> str:
    pdf = PdfReader(str(file))
    extracted = [page.extract_text() or "" for page in pdf.pages]
    return "\n".join(extracted).strip()

def extract_pptx(file: Path) -> str:
    pptx = Presentation(str(file))
    extracted = []
    for slide in pptx.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                extracted.append(shape.text)
    return "\n".join(extracted).strip() 

def get_name(file_name: str) -> str:
    name = re.sub(r"[^A-Za-z0-9._-]+", "_", file_name).strip("._")
    return name or "slides"

async def save_file(file: UploadFile) -> Path: 
    file_type = Path(file.filename or "").suffix.lower()
    file_name = get_name(Path(file.filename or "slides").stem)
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    saved_name = f"{file_name}_{timestamp}{file_type}" 

    UPLOAD_DIR.mkdir(parents=True, exist_ok=True)
    dest = UPLOAD_DIR / saved_name

    size = 0
    with dest.open("wb") as f:
        while True:
            chunk = await file.read(1024 * 1024)
            if not chunk:
                break
            size += len(chunk)
            if size > MAX_UPLOAD_SIZE_BYTES:
                dest.unlink(missing_ok=True)
                raise HTTPException(
                    status_code=413,
                    detail="File is too large. Please upload a file under 25MB.",
                )
            f.write(chunk)
    await file.close()
    return dest

def chunk_text(content: str, chunk_size: int = 12000, overlap: int = 1000) -> list[str]:
    "Split long extracted text into chunks"
    cleaned = content.strip()
    if not cleaned: 
        return []
    if len(cleaned) <= chunk_size:
        return [cleaned]

    chunks = []
    step = max(1, chunk_size - overlap)
    start = 0
    while start < len(cleaned):
        end = min(len(cleaned), start + chunk_size)
        chunk = cleaned[start: end].strip()
        if chunk: chunks.append(chunk)
        if end >= len(cleaned):
            break
        start += step
    return chunks


async def summarize_chunk(index: int, chunk: str, chunks: list[str]) -> tuple[list[str], list[str]]:
    chunk_prompt = build_chunk_summary_prompt(chunk, index + 1, len(chunks))
    chunk_response = await gemini_service.call_gemini(chunk_prompt) 
    if chunk_response is None:
            raise HTTPException(status_code=500, detail="Failed to summarize chunk of slide.")
    try:
        cleaned = json.loads(clean_response(chunk_response))
        response_summary = cleaned.get("summary",[])
        response_facts = cleaned.get("key_facts", [])
        summary = []
        key_facts = []
        if isinstance(response_summary, list):
            summary = [str(s).strip() for s in response_summary if str(s).strip()]
        if isinstance(response_facts, list):
            key_facts = [str(kf).strip() for kf in response_facts if str(kf).strip()]
        return summary, key_facts
    except json.JSONDecodeError as e:
            logger.warning("[slides chunk] Failed to parse Gemini JSON: %s", e)
            raise HTTPException(
                status_code=500,
                detail="Failed to parse AI response as JSON. Please try again."
            )    

async def get_study_pack_responses(text: str):
    quiz_prompt = build_quiz_generation_prompt(content=text)
    flashcards_prompt = build_flashcard_generation_prompt(content=text, mode="notes")
    quiz_response, flashcards_response = await asyncio.gather(
        gemini_service.call_gemini(quiz_prompt),
        gemini_service.call_gemini(flashcards_prompt),
    )
    if quiz_response is None:
        raise HTTPException(status_code=500, detail="Failed to generate quiz.")
    if flashcards_response is None:
        raise HTTPException(status_code=500, detail="Failed to generate flashcards.")
    return quiz_response, flashcards_response
    

async def generate_study_pack_slides(text: str) -> tuple[list[str], list[QuizQuestion], list[Flashcard]]:
    text = text.strip()
    # if no text
    if not text: 
        raise HTTPException(status_code=422, detail="No text found in slides, not able to generate study pack")
    # if text not require chunking
    if len(text) < CHUNKING_THRESHOLD_CHARS:
        study_prompt = build_summary_only_prompt(user_notes=text, num_points=5)
        study_response = await gemini_service.call_gemini(study_prompt)
        if study_response is None:
            raise HTTPException(status_code=500, detail="Failed to generate summary.")
        
        quiz_response, flashcards_response = await get_study_pack_responses(text)

        try:
            study_data = json.loads(clean_response(study_response))
            summary = study_data.get("summary", [])
            if not isinstance(summary, list):
                raise ValueError("Response missing 'summary' array")
            summary = [str(item).strip() for item in summary if str(item).strip()]
            if not summary:
                raise ValueError("Response missing non-empty summary items")
            quiz_questions = parse_and_validate_quiz(quiz_response)
            flashcards = parse_and_validate_flashcards(flashcards_response)
            return summary, quiz_questions, flashcards
        except json.JSONDecodeError as e:
            logger.warning("[slides] Failed to parse Gemini JSON: %s", e)
            raise HTTPException(
                status_code=500,
                detail="Failed to parse AI response as JSON. Please try again."
            )
        except (KeyError, TypeError, ValueError) as e:
            logger.warning("[slides] Invalid Gemini response structure: %s", e)
            raise HTTPException(
                status_code=500,
                detail=f"Invalid AI response format: {str(e)}"
            )

    # else, chunk the text and merge
    chunks = chunk_text(text)
    if not chunks:
        raise HTTPException(status_code=422, detail="No text found in slides, not able to generate study pack")

    chunk_results = await asyncio.gather(
        *(summarize_chunk(index, chunk, chunks) for index, chunk in enumerate(chunks))
    )

    chunk_summaries = []
    chunk_facts = []
    for summaries, facts in chunk_results:
        chunk_summaries.extend(summaries)
        chunk_facts.extend(facts)
    
    if not chunk_summaries and not chunk_facts:
        raise HTTPException(status_code=500, detail="No usable content produced from slide chunks.")
    
    merge_prompt = build_merge_outline_prompt(chunk_summaries, chunk_facts)
    merge_response = await gemini_service.call_gemini(merge_prompt)
    if not merge_response:
        raise HTTPException(status_code=422, detail="Failed to merge chunks summaries/facts.")
    
    try:
        merged = json.loads(clean_response(merge_response))
        outline = merged.get("outline", [])
        summary = merged.get("summary", [])
        if not isinstance(outline, list) or not outline:
            raise ValueError("Merged outline missing or empty")
        if not isinstance(summary, list) or not summary:
            raise ValueError("Merged summary missing or empty")
        canonical_outline = "\n".join(f"- {str(item).strip()}" for item in outline if str(item).strip())
        if not canonical_outline.strip():
            raise ValueError("Merged outline empty after cleanup")
        summary = [str(item).strip() for item in summary if str(item).strip()]
        if not summary:
            raise ValueError("Merged summary empty after cleanup")
    except json.JSONDecodeError as e:
            logger.warning("[slides merge] Failed to parse Gemini JSON: %s", e)
            raise HTTPException(
                status_code=500,
                detail="Failed to parse AI response as JSON. Please try again."
            )
    
    quiz_response, flashcards_response = await get_study_pack_responses(canonical_outline)
    try:
        quiz_questions = parse_and_validate_quiz(quiz_response)
        flashcards = parse_and_validate_flashcards(flashcards_response)
        return summary, quiz_questions, flashcards
    except json.JSONDecodeError as e:
        logger.warning("[slides merge] Failed to parse Gemini JSON: %s", e)
        raise HTTPException(
            status_code=500,
            detail="Failed to parse AI response as JSON. Please try again."
        )
    except (KeyError, TypeError, ValueError) as e:
        logger.warning("[slides merge] Invalid Gemini response structure: %s", e)
        raise HTTPException(
            status_code=500,
            detail=f"Invalid AI response format: {str(e)}"
        )

@app.post("/api/v1/slides/study-pack", response_model=SlideStudyPackResponse)
async def generate_study_pack_from_slides(
        file: UploadFile = File(...), 
        user: Optional[UserPayload] = Depends(user_for_generate),
):
    filename = file.filename or "slides"
    suffix = Path(filename).suffix.lower()
    valid_suffixes = {".pdf", ".pptx"}
    valid_content_types = {
        "application/pdf",
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    }

    if suffix not in valid_suffixes and (file.content_type or "") not in valid_content_types:
        raise HTTPException(
            status_code=400,
            detail="Unsupported file type. Please upload a PDF or PPTX file.",
        )

    stored_path = await save_file(file)

    try:
        if stored_path.suffix.lower() == ".pdf":
            extracted = extract_pdf(stored_path)
        else:
            extracted = extract_pptx(stored_path)
    except Exception as e:
        logger.warning("[slide extract] Slide extraction failed for %s: %s", stored_path, e)
        raise HTTPException(
            status_code=400,
            detail="Could not extract text from this file. Please try another PDF or PPTX.",
        )
    
    if len(extracted.strip()) == 0:
        raise HTTPException(
            status_code=422,
            detail="No text found in uploaded slides, cannot generate study pack.",
        )
    
    summary, quiz_questions, flashcards = await generate_study_pack_slides(extracted)
    quiz_set_id = None
    flashcard_set_id = None
    # save quiz and flashcards to database
    sb = get_supabase()
    try:
        quiz_result = sb.table("quiz").insert({
            "user_id": user["user_id"] if user else "00000000-0000-0000-0000-000000000001",
            "source_text": f"Generated from slides: {filename}",
            "questions": [q.model_dump() for q in quiz_questions],
        }).execute()
        quiz_set_id = quiz_result.data[0]["id"]
    except Exception as e:
        logger.warning("[slides study-pack] Failed to store quiz: %s", e)
        raise HTTPException(
            status_code=500,
            detail="Failed to store generated quiz. Please try again.",
        )
    try:
        flashcards_result = sb.table("flashcards").insert({
            "user_id": user["user_id"] if user else "00000000-0000-0000-0000-000000000001",
            "source_text": f"Generated from slides: {filename}",
            "topic": "Slides upload",
            "cards": [fc.model_dump() for fc in flashcards],
        }).execute()
        flashcard_set_id = flashcards_result.data[0]["id"]
    except Exception as e:
        logger.warning("[slides study-pack] Failed to store flashcards: %s", e)
        raise HTTPException(
            status_code=500,
            detail="Failed to store generated flashcards. Please try again.",
        )

    return SlideStudyPackResponse(
        file_name=filename,
        stored_path=str(stored_path),
        quiz_set_id=quiz_set_id,
        flashcard_set_id=flashcard_set_id,
        extracted_text=extracted,
        summary=summary,
        quiz=quiz_questions,
        flashcards=flashcards,
    )


@app.post("/api/v1/slides/quiz/regenerate", response_model=GenerateQuizResponse)
async def regenerate_slide_quiz_questions(
    request: SlideQuizRegenerateRequest,
    user: UserPayload | None = Depends(user_for_generate),
):
    """Regenerate quiz from extracted slide text without StudyPackRequest size cap."""
    prompt = build_quiz_generation_prompt(content=request.text)
    response = await gemini_service.call_gemini(prompt)

    if response is None:
        raise HTTPException(
            status_code=500,
            detail="Failed to regenerate quiz. Please try again.",
        )

    quiz_questions = parse_and_validate_quiz(response)

    sb = get_supabase()
    try:
        result = sb.table("quiz").insert({
            "user_id": user["user_id"] if user else "00000000-0000-0000-0000-000000000001",
            "source_text": request.text,
            "questions": [q.model_dump() for q in quiz_questions],
        }).execute()
        quiz_set_id = result.data[0]["id"]
    except Exception as e:
        logger.warning("[slides quiz regenerate] DB insert failed: %s", e)
        raise HTTPException(
            status_code=500,
            detail="Failed to store regenerated quiz. Please try again.",
        )

    return GenerateQuizResponse(quiz_set_id=quiz_set_id, quiz=quiz_questions)
    